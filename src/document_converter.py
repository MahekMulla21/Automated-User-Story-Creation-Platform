import os


def convert_document_to_markdown(file_path: str) -> str:
    """
    Dispatches to the right extractor based on file extension and
    returns the extracted content as a markdown-ish string.

    Raises:
        ValueError: unsupported file extension.
        RuntimeError: no extractable text was found in the file.
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return _convert_pdf(file_path)
    elif ext == ".docx":
        return _convert_docx(file_path)
    elif ext == ".pptx":
        return _convert_pptx(file_path)
    elif ext == ".xlsx":
        return _convert_xlsx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def _convert_pdf(file_path: str) -> str:
    import pdfplumber

    text_parts = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_parts.append(f"### Page {i}\n{page_text.strip()}")

    result = "\n\n".join(text_parts).strip()
    if not result:
        raise RuntimeError(
            "No extractable text found in the PDF (it may be a scanned "
            "image with no text layer)."
        )
    return result


def _convert_docx(file_path: str) -> str:
    import docx

    document = docx.Document(file_path)
    lines = []

    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower()
        if style.startswith("heading"):
            digits = "".join(ch for ch in style if ch.isdigit())
            level = int(digits) if digits else 1
            level = max(1, min(level, 6))
            lines.append(f"{'#' * level} {text}")
        else:
            lines.append(text)

    for table in document.tables:
        lines.append("")
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            lines.append(" | ".join(cells))

    result = "\n".join(lines).strip()
    if not result:
        raise RuntimeError("No extractable text found in the DOCX file.")
    return result


def _convert_pptx(file_path: str) -> str:
    from pptx import Presentation

    prs = Presentation(file_path)
    lines = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        slide_lines.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_lines.append(" | ".join(cells))

        if slide_lines:
            lines.append(f"### Slide {i}")
            lines.extend(slide_lines)
            lines.append("")

    result = "\n".join(lines).strip()
    if not result:
        raise RuntimeError("No extractable text found in the PPTX file.")
    return result


def _convert_xlsx(file_path: str) -> str:
    import openpyxl

    workbook = openpyxl.load_workbook(file_path, data_only=True)
    lines = []

    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        sheet_lines = []
        for row in sheet.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                cells = ["" if c is None else str(c) for c in row]
                sheet_lines.append(" | ".join(cells))

        if sheet_lines:
            lines.append(f"### Sheet: {sheet_name}")
            lines.extend(sheet_lines)
            lines.append("")

    result = "\n".join(lines).strip()
    if not result:
        raise RuntimeError("No extractable text found in the XLSX file.")
    return result
